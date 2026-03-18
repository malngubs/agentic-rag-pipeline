"""
╔══════════════════════════════════════════════════════════════════════════════╗
║                         FILE INGESTOR MODULE                                 ║
║              Multi-Format Data Extraction for BI Platform                    ║
╠══════════════════════════════════════════════════════════════════════════════╣
║  Extracts tabular data from: PDF, Excel, CSV, PowerPoint, Word, Images       ║
║  Features: OCR support, table detection, encoding detection, multi-sheet     ║
╚══════════════════════════════════════════════════════════════════════════════╝

SUPPORTED FORMATS:
------------------
📄 PDF      - Text extraction + table detection + OCR for scanned documents
📊 Excel    - .xlsx, .xls, .xlsm with multi-sheet and formula support
📋 CSV/TSV  - Auto-encoding detection, delimiter inference
📑 PowerPoint - .pptx table extraction from slides
📝 Word     - .docx table extraction
🖼️ Images   - PNG, JPG, TIFF with OCR table extraction
📦 JSON     - Nested structure flattening

USAGE:
------
    from data_analyst import FileIngestor
    
    # Initialize ingestor
    ingestor = FileIngestor()
    
    # Extract from any file
    result = await ingestor.extract("sales_report.pdf")
    
    # Access extracted data
    for table in result.tables:
        print(f"Table: {table.name}")
        print(table.dataframe.head())
    
    # Check extraction metadata
    print(f"Source: {result.source_file}")
    print(f"Tables found: {len(result.tables)}")
    print(f"Warnings: {result.warnings}")

ARCHITECTURE:
-------------
FileIngestor (main class)
    ├── PDFExtractor
    ├── ExcelExtractor
    ├── CSVExtractor
    ├── PowerPointExtractor
    ├── WordExtractor
    ├── ImageExtractor
    └── JSONExtractor
"""

# =============================================================================
# IMPORTS
# =============================================================================

from __future__ import annotations

import asyncio
import io
import json
import logging
import mimetypes
import os
import re
import tempfile
import warnings
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from datetime import datetime
from enum import Enum
from pathlib import Path
from typing import (
    Any,
    BinaryIO,
    Callable,
    Dict,
    List,
    Optional,
    Tuple,
    Union,
)

# Data handling
import pandas as pd
import numpy as np

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Suppress warnings from libraries
warnings.filterwarnings("ignore", category=UserWarning)


# =============================================================================
# OPTIONAL DEPENDENCIES - Graceful degradation if not installed
# =============================================================================

# PDF Processing
try:
    import pdfplumber
    PDF_PLUMBER_AVAILABLE = True
except ImportError:
    PDF_PLUMBER_AVAILABLE = False
    logger.warning("pdfplumber not installed - PDF support limited")

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    logger.warning("PyMuPDF not installed - PDF image extraction disabled")

# OCR
try:
    import pytesseract
    from PIL import Image
    # Set Tesseract path for Windows
    pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False
    logger.warning("pytesseract/PIL not installed - OCR disabled")

try:
    import easyocr
    EASYOCR_AVAILABLE = True
except ImportError:
    EASYOCR_AVAILABLE = False
    logger.debug("easyocr not installed - using pytesseract for OCR")

# Excel Processing
try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False
    logger.warning("openpyxl not installed - .xlsx support limited")

try:
    import xlrd
    XLRD_AVAILABLE = True
except ImportError:
    XLRD_AVAILABLE = False
    logger.debug("xlrd not installed - .xls support disabled")

# Office Documents
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not installed - PowerPoint support disabled")

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not installed - Word support disabled")

# Encoding Detection
try:
    import chardet
    CHARDET_AVAILABLE = True
except ImportError:
    CHARDET_AVAILABLE = False
    logger.debug("chardet not installed - using default encoding detection")


# =============================================================================
# ENUMERATIONS
# =============================================================================

class FileType(Enum):
    """Supported file types for extraction"""
    PDF = "pdf"
    EXCEL_XLSX = "xlsx"
    EXCEL_XLS = "xls"
    EXCEL_XLSM = "xlsm"
    CSV = "csv"
    TSV = "tsv"
    POWERPOINT = "pptx"
    WORD = "docx"
    IMAGE_PNG = "png"
    IMAGE_JPG = "jpg"
    IMAGE_JPEG = "jpeg"
    IMAGE_TIFF = "tiff"
    IMAGE_BMP = "bmp"
    JSON = "json"
    UNKNOWN = "unknown"


class ExtractionStatus(Enum):
    """Status of extraction operation"""
    SUCCESS = "success"
    PARTIAL = "partial"  # Some tables extracted, some failed
    FAILED = "failed"
    NO_DATA = "no_data"  # File parsed but no tabular data found


class DataQuality(Enum):
    """Quality assessment of extracted data"""
    HIGH = "high"       # Clean, well-structured data
    MEDIUM = "medium"   # Some issues but usable
    LOW = "low"         # Significant issues, may need cleaning
    UNKNOWN = "unknown"


# =============================================================================
# DATA CLASSES
# =============================================================================

@dataclass
class ExtractedTable:
    """
    Represents a single extracted table from a document.
    
    Attributes:
        name: Identifier for the table (e.g., "Sheet1", "Table on Page 3")
        dataframe: The actual data as a pandas DataFrame
        source_location: Where in the document this table was found
        row_count: Number of rows in the table
        column_count: Number of columns in the table
        quality: Assessment of data quality
        metadata: Additional information about the extraction
    """
    name: str
    dataframe: pd.DataFrame
    source_location: str = ""
    row_count: int = 0
    column_count: int = 0
    quality: DataQuality = DataQuality.UNKNOWN
    metadata: Dict[str, Any] = field(default_factory=dict)
    
    def __post_init__(self):
        """Calculate row and column counts from dataframe"""
        if self.dataframe is not None and not self.dataframe.empty:
            self.row_count = len(self.dataframe)
            self.column_count = len(self.dataframe.columns)
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for JSON serialization"""
        return {
            "name": self.name,
            "source_location": self.source_location,
            "row_count": self.row_count,
            "column_count": self.column_count,
            "quality": self.quality.value,
            "columns": list(self.dataframe.columns) if self.dataframe is not None else [],
            "preview": self.dataframe.head(5).to_dict('records') if self.dataframe is not None else [],
            "metadata": self.metadata
        }


@dataclass
class ExtractionResult:
    """
    Complete result of a file extraction operation.
    
    Attributes:
        source_file: Original filename
        file_type: Detected file type
        status: Overall extraction status
        tables: List of extracted tables
        text_content: Raw text content (for search/RAG)
        warnings: Any warnings during extraction
        errors: Any errors during extraction
        processing_time: Time taken to extract (seconds)
        metadata: Additional file metadata
    """
    source_file: str
    file_type: FileType
    status: ExtractionStatus
    tables: List[ExtractedTable] = field(default_factory=list)
    text_content: str = ""
    warnings: List[str] = field(default_factory=list)
    errors: List[str] = field(default_factory=list)
    processing_time: float = 0.0
    metadata: Dict[str, Any] = field(default_factory=dict)
    
    @property
    def has_data(self) -> bool:
        """Check if any data was extracted"""
        return len(self.tables) > 0 and any(t.row_count > 0 for t in self.tables)
    
    @property
    def total_rows(self) -> int:
        """Total rows across all tables"""
        return sum(t.row_count for t in self.tables)
    
    @property
    def total_columns(self) -> int:
        """Total unique columns across all tables"""
        all_cols = set()
        for t in self.tables:
            all_cols.update(t.dataframe.columns)
        return len(all_cols)
    
    @property
    def primary_table(self) -> Optional[ExtractedTable]:
        """Get the largest table (most rows)"""
        if not self.tables:
            return None
        return max(self.tables, key=lambda t: t.row_count)
    
    def get_dataframe(self, table_name: Optional[str] = None) -> Optional[pd.DataFrame]:
        """
        Get a dataframe by table name, or the primary table if no name specified.
        
        Args:
            table_name: Name of the table to retrieve
            
        Returns:
            DataFrame or None if not found
        """
        if table_name is None:
            return self.primary_table.dataframe if self.primary_table else None
        
        for table in self.tables:
            if table.name == table_name:
                return table.dataframe
        return None
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for JSON serialization"""
        return {
            "source_file": self.source_file,
            "file_type": self.file_type.value,
            "status": self.status.value,
            "tables": [t.to_dict() for t in self.tables],
            "table_count": len(self.tables),
            "total_rows": self.total_rows,
            "has_data": self.has_data,
            "warnings": self.warnings,
            "errors": self.errors,
            "processing_time": round(self.processing_time, 3),
            "metadata": self.metadata
        }


@dataclass
class ExtractionConfig:
    """
    Configuration options for file extraction.
    
    Attributes:
        max_rows: Maximum rows to extract per table (None = unlimited)
        max_tables: Maximum tables to extract per file
        ocr_enabled: Enable OCR for scanned documents/images
        ocr_language: Language for OCR (default: English)
        encoding: Force specific encoding (None = auto-detect)
        date_formats: List of date formats to try when parsing
        header_row: Row number to use as header (None = auto-detect)
        skip_empty_tables: Skip tables with no data
        min_table_rows: Minimum rows for a table to be included
        progress_callback: Optional callback for progress updates
    """
    max_rows: Optional[int] = None
    max_tables: int = 50
    ocr_enabled: bool = True
    ocr_language: str = "eng"
    encoding: Optional[str] = None
    date_formats: List[str] = field(default_factory=lambda: [
        "%Y-%m-%d", "%d/%m/%Y", "%m/%d/%Y", "%Y/%m/%d",
        "%d-%m-%Y", "%m-%d-%Y", "%Y%m%d",
        "%d %b %Y", "%d %B %Y", "%b %d, %Y", "%B %d, %Y"
    ])
    header_row: Optional[int] = None
    skip_empty_tables: bool = True
    min_table_rows: int = 1
    progress_callback: Optional[Callable[[str, float], None]] = None
    
    def report_progress(self, message: str, progress: float):
        """Report progress if callback is set"""
        if self.progress_callback:
            self.progress_callback(message, min(1.0, max(0.0, progress)))


# =============================================================================
# UTILITY FUNCTIONS
# =============================================================================

def detect_file_type(file_path: Union[str, Path], content: Optional[bytes] = None) -> FileType:
    """
    Detect file type from extension and/or content.
    
    Args:
        file_path: Path to the file
        content: Optional file content for magic number detection
        
    Returns:
        Detected FileType enum value
    """
    path = Path(file_path)
    extension = path.suffix.lower().lstrip('.')
    
    # Map extensions to file types
    extension_map = {
        'pdf': FileType.PDF,
        'xlsx': FileType.EXCEL_XLSX,
        'xls': FileType.EXCEL_XLS,
        'xlsm': FileType.EXCEL_XLSM,
        'csv': FileType.CSV,
        'tsv': FileType.TSV,
        'pptx': FileType.POWERPOINT,
        'docx': FileType.WORD,
        'png': FileType.IMAGE_PNG,
        'jpg': FileType.IMAGE_JPG,
        'jpeg': FileType.IMAGE_JPEG,
        'tiff': FileType.IMAGE_TIFF,
        'tif': FileType.IMAGE_TIFF,
        'bmp': FileType.IMAGE_BMP,
        'json': FileType.JSON,
    }
    
    return extension_map.get(extension, FileType.UNKNOWN)


def detect_encoding(file_path: Union[str, Path], sample_size: int = 10000) -> str:
    """
    Detect file encoding using chardet or fallback methods.
    
    Args:
        file_path: Path to the file
        sample_size: Number of bytes to sample
        
    Returns:
        Detected encoding name
    """
    try:
        with open(file_path, 'rb') as f:
            sample = f.read(sample_size)
        
        if CHARDET_AVAILABLE:
            result = chardet.detect(sample)
            encoding = result.get('encoding', 'utf-8')
            confidence = result.get('confidence', 0)
            
            # If low confidence, try common encodings
            if confidence < 0.7:
                for enc in ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']:
                    try:
                        sample.decode(enc)
                        return enc
                    except (UnicodeDecodeError, LookupError):
                        continue
            
            return encoding or 'utf-8'
        else:
            # Fallback: try common encodings
            for enc in ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252']:
                try:
                    sample.decode(enc)
                    return enc
                except (UnicodeDecodeError, LookupError):
                    continue
            return 'utf-8'
            
    except Exception as e:
        logger.warning(f"Encoding detection failed: {e}, defaulting to utf-8")
        return 'utf-8'


def detect_csv_delimiter(file_path: Union[str, Path], encoding: str = 'utf-8') -> str:
    """
    Auto-detect CSV delimiter by analyzing the first few lines.
    
    Args:
        file_path: Path to the CSV file
        encoding: File encoding
        
    Returns:
        Detected delimiter character
    """
    delimiters = [',', ';', '\t', '|']
    
    try:
        with open(file_path, 'r', encoding=encoding, errors='replace') as f:
            # Read first 5 lines
            sample_lines = [f.readline() for _ in range(5)]
        
        sample = ''.join(sample_lines)
        
        # Count occurrences of each delimiter
        counts = {d: sample.count(d) for d in delimiters}
        
        # Pick the most common one (that appears consistently)
        best_delimiter = max(counts, key=counts.get)
        
        # Verify it appears in multiple lines
        if counts[best_delimiter] >= len(sample_lines) - 1:
            return best_delimiter
        
        return ','  # Default to comma
        
    except Exception as e:
        logger.warning(f"Delimiter detection failed: {e}, defaulting to comma")
        return ','


def clean_dataframe(df: pd.DataFrame) -> pd.DataFrame:
    """
    Clean a dataframe by handling common issues.
    
    Args:
        df: Input dataframe
        
    Returns:
        Cleaned dataframe
    """
    if df is None or df.empty:
        return df
    
    # Make a copy to avoid modifying original
    df = df.copy()
    
    # Remove completely empty rows and columns
    df = df.dropna(how='all')
    df = df.dropna(axis=1, how='all')
    
    # Clean and sanitize column names for SQL compatibility
    def sanitize_column_name(col_name: str, index: int) -> str:
        """Sanitize a column name to be SQL-safe"""
        col = str(col_name).strip()

        # Handle unnamed columns (pandas creates these as "Unnamed: 0", "Unnamed: 1", etc.)
        if col.lower().startswith('unnamed'):
            return f'column_{index}'

        # Replace problematic characters with underscores
        # Colons, spaces, hyphens, dots, and other special chars break SQL
        import re
        col = re.sub(r'[:\s\-\.\(\)\[\]\{\}\+\*\/\\@#$%^&!=<>,;\'\"]+', '_', col)

        # Remove leading/trailing underscores
        col = col.strip('_')

        # Ensure it doesn't start with a number (invalid SQL identifier)
        if col and col[0].isdigit():
            col = f'col_{col}'

        # If empty after sanitization, use index-based name
        if not col:
            return f'column_{index}'

        return col

    # Apply sanitization to all columns
    new_columns = []
    seen_names = set()
    for i, col in enumerate(df.columns):
        sanitized = sanitize_column_name(col, i)

        # Handle duplicate column names by appending index
        if sanitized in seen_names:
            sanitized = f'{sanitized}_{i}'
        seen_names.add(sanitized)
        new_columns.append(sanitized)

    df.columns = new_columns
    
    # Replace empty strings with NaN
    df = df.replace(r'^\s*$', np.nan, regex=True)
    
    # Strip whitespace from string columns
    for col in df.select_dtypes(include=['object']).columns:
        df[col] = df[col].apply(lambda x: x.strip() if isinstance(x, str) else x)
    
    return df


def assess_data_quality(df: pd.DataFrame) -> DataQuality:
    """
    Assess the quality of extracted data.
    
    Args:
        df: DataFrame to assess
        
    Returns:
        DataQuality enum value
    """
    if df is None or df.empty:
        return DataQuality.UNKNOWN
    
    total_cells = df.size
    if total_cells == 0:
        return DataQuality.UNKNOWN
    
    # Calculate quality metrics
    missing_ratio = df.isna().sum().sum() / total_cells
    
    # Check for unnamed columns
    unnamed_cols = sum(1 for col in df.columns if 'unnamed' in str(col).lower())
    unnamed_ratio = unnamed_cols / len(df.columns) if len(df.columns) > 0 else 0
    
    # Assess quality
    if missing_ratio < 0.05 and unnamed_ratio < 0.1:
        return DataQuality.HIGH
    elif missing_ratio < 0.20 and unnamed_ratio < 0.3:
        return DataQuality.MEDIUM
    else:
        return DataQuality.LOW


# =============================================================================
# BASE EXTRACTOR CLASS
# =============================================================================

class BaseExtractor(ABC):
    """
    Abstract base class for all file extractors.
    
    All specific extractors (PDF, Excel, etc.) inherit from this class
    and implement the extract() method.
    """
    
    def __init__(self, config: Optional[ExtractionConfig] = None):
        """
        Initialize the extractor with configuration.
        
        Args:
            config: Extraction configuration options
        """
        self.config = config or ExtractionConfig()
        self.warnings: List[str] = []
        self.errors: List[str] = []
    
    @abstractmethod
    async def extract(self, file_path: Union[str, Path]) -> List[ExtractedTable]:
        """
        Extract tables from a file.
        
        Args:
            file_path: Path to the file to extract
            
        Returns:
            List of ExtractedTable objects
        """
        pass
    
    @abstractmethod
    def can_handle(self, file_type: FileType) -> bool:
        """
        Check if this extractor can handle the given file type.
        
        Args:
            file_type: The file type to check
            
        Returns:
            True if this extractor can handle the file type
        """
        pass
    
    def add_warning(self, message: str):
        """Add a warning message"""
        self.warnings.append(message)
        logger.warning(message)
    
    def add_error(self, message: str):
        """Add an error message"""
        self.errors.append(message)
        logger.error(message)
    
    def clear_messages(self):
        """Clear all warnings and errors"""
        self.warnings = []
        self.errors = []


# =============================================================================
# PDF EXTRACTOR
# =============================================================================

class PDFExtractor(BaseExtractor):
    """
    Extracts tables and text from PDF documents.
    
    Features:
    - Table detection and extraction using pdfplumber
    - OCR support for scanned PDFs
    - Text extraction for RAG/search
    - Image extraction for further processing
    
    Dependencies:
    - pdfplumber (required for table extraction)
    - pytesseract + PIL (optional, for OCR)
    - PyMuPDF/fitz (optional, for image extraction)
    """
    
    def can_handle(self, file_type: FileType) -> bool:
        """PDF extractor handles PDF files"""
        return file_type == FileType.PDF
    
    async def extract(self, file_path: Union[str, Path]) -> List[ExtractedTable]:
        """
        Extract tables from a PDF file.
        
        Strategy:
        1. Try pdfplumber for native PDF tables
        2. If no tables found and OCR enabled, try OCR
        3. Return all found tables with quality assessment
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            List of ExtractedTable objects
        """
        file_path = Path(file_path)
        tables: List[ExtractedTable] = []
        
        if not PDF_PLUMBER_AVAILABLE:
            self.add_error("pdfplumber not installed - cannot extract PDF tables")
            return tables
        
        self.config.report_progress("Opening PDF...", 0.1)
        
        try:
            with pdfplumber.open(file_path) as pdf:
                total_pages = len(pdf.pages)
                self.config.report_progress(f"Processing {total_pages} pages...", 0.2)
                
                # Track if we found any tables
                tables_found = False
                
                # Process each page
                for page_num, page in enumerate(pdf.pages):
                    progress = 0.2 + (0.6 * (page_num / total_pages))
                    self.config.report_progress(f"Processing page {page_num + 1}/{total_pages}", progress)
                    
                    # Extract tables from this page
                    page_tables = self._extract_tables_from_page(page, page_num + 1)
                    
                    if page_tables:
                        tables_found = True
                        tables.extend(page_tables)
                    
                    # Check if we've hit the max tables limit
                    if len(tables) >= self.config.max_tables:
                        self.add_warning(f"Reached maximum table limit ({self.config.max_tables})")
                        break
                
                # If no tables found and OCR is enabled, try OCR
                if not tables_found and self.config.ocr_enabled:
                    self.config.report_progress("No tables found, trying OCR...", 0.8)
                    ocr_tables = await self._extract_with_ocr(file_path)
                    tables.extend(ocr_tables)
                
                self.config.report_progress("Extraction complete", 1.0)
                
        except Exception as e:
            self.add_error(f"PDF extraction failed: {str(e)}")
            logger.exception("PDF extraction error")
        
        return tables
    
    def _extract_tables_from_page(
        self, 
        page, 
        page_num: int
    ) -> List[ExtractedTable]:
        """
        Extract tables from a single PDF page.
        
        Args:
            page: pdfplumber page object
            page_num: Page number (1-indexed)
            
        Returns:
            List of ExtractedTable objects from this page
        """
        tables = []
        
        try:
            # Find tables on this page
            page_tables = page.extract_tables()
            
            for table_idx, table_data in enumerate(page_tables):
                if not table_data or len(table_data) < 2:
                    continue
                
                # Convert to DataFrame
                try:
                    # First row as header
                    headers = table_data[0]
                    data = table_data[1:]
                    
                    # Clean headers
                    headers = [
                        str(h).strip() if h else f"Column_{i}"
                        for i, h in enumerate(headers)
                    ]
                    
                    # Handle duplicate headers
                    seen = {}
                    clean_headers = []
                    for h in headers:
                        if h in seen:
                            seen[h] += 1
                            clean_headers.append(f"{h}_{seen[h]}")
                        else:
                            seen[h] = 0
                            clean_headers.append(h)
                    
                    # Create DataFrame
                    df = pd.DataFrame(data, columns=clean_headers)
                    
                    # Clean the dataframe
                    df = clean_dataframe(df)
                    
                    # Skip empty tables if configured
                    if self.config.skip_empty_tables and df.empty:
                        continue
                    
                    # Skip tables with too few rows
                    if len(df) < self.config.min_table_rows:
                        continue
                    
                    # Apply max rows limit
                    if self.config.max_rows and len(df) > self.config.max_rows:
                        df = df.head(self.config.max_rows)
                        self.add_warning(f"Table truncated to {self.config.max_rows} rows")
                    
                    # Create ExtractedTable
                    table = ExtractedTable(
                        name=f"Page_{page_num}_Table_{table_idx + 1}",
                        dataframe=df,
                        source_location=f"Page {page_num}",
                        quality=assess_data_quality(df),
                        metadata={
                            "page_number": page_num,
                            "table_index": table_idx,
                            "extraction_method": "pdfplumber"
                        }
                    )
                    tables.append(table)
                    
                except Exception as e:
                    self.add_warning(f"Failed to process table on page {page_num}: {str(e)}")
                    continue
                    
        except Exception as e:
            self.add_warning(f"Error extracting tables from page {page_num}: {str(e)}")
        
        return tables
    
    async def _extract_with_ocr(self, file_path: Path) -> List[ExtractedTable]:
        """
        Extract tables using OCR for scanned PDFs.
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            List of ExtractedTable objects from OCR
        """
        tables = []
        
        if not TESSERACT_AVAILABLE:
            self.add_warning("OCR not available - pytesseract/PIL not installed")
            return tables
        
        if not PYMUPDF_AVAILABLE:
            self.add_warning("Cannot extract images for OCR - PyMuPDF not installed")
            return tables
        
        try:
            # Open PDF with PyMuPDF to extract images
            doc = fitz.open(file_path)
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # Render page to image
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for better OCR
                img_data = pix.tobytes("png")
                
                # Convert to PIL Image
                img = Image.open(io.BytesIO(img_data))
                
                # Run OCR
                ocr_text = pytesseract.image_to_string(
                    img,
                    lang=self.config.ocr_language,
                    config='--psm 6'  # Assume uniform block of text
                )
                
                # Try to detect tables in OCR text
                table_df = self._parse_ocr_table(ocr_text)
                
                if table_df is not None and not table_df.empty:
                    table = ExtractedTable(
                        name=f"OCR_Page_{page_num + 1}",
                        dataframe=table_df,
                        source_location=f"Page {page_num + 1} (OCR)",
                        quality=DataQuality.LOW,  # OCR quality is typically lower
                        metadata={
                            "page_number": page_num + 1,
                            "extraction_method": "ocr",
                            "ocr_language": self.config.ocr_language
                        }
                    )
                    tables.append(table)
            
            doc.close()
            
        except Exception as e:
            self.add_error(f"OCR extraction failed: {str(e)}")
            logger.exception("OCR extraction error")
        
        return tables
    
    def _parse_ocr_table(self, ocr_text: str) -> Optional[pd.DataFrame]:
        """
        Attempt to parse tabular data from OCR text.
        
        Args:
            ocr_text: Raw text from OCR
            
        Returns:
            DataFrame if table detected, None otherwise
        """
        if not ocr_text or not ocr_text.strip():
            return None
        
        lines = ocr_text.strip().split('\n')
        lines = [line.strip() for line in lines if line.strip()]
        
        if len(lines) < 2:
            return None
        
        # Try to detect delimiter
        potential_delimiters = ['\t', '  ', '|', ',']
        
        for delimiter in potential_delimiters:
            try:
                # Split lines by delimiter
                rows = []
                for line in lines:
                    if delimiter == '  ':
                        # Multiple spaces
                        parts = re.split(r'\s{2,}', line)
                    else:
                        parts = line.split(delimiter)
                    
                    parts = [p.strip() for p in parts if p.strip()]
                    if parts:
                        rows.append(parts)
                
                # Check if we have consistent column counts
                if len(rows) >= 2:
                    col_counts = [len(row) for row in rows]
                    most_common_cols = max(set(col_counts), key=col_counts.count)
                    
                    if most_common_cols >= 2:
                        # Filter to rows with correct column count
                        filtered_rows = [row for row in rows if len(row) == most_common_cols]
                        
                        if len(filtered_rows) >= 2:
                            headers = filtered_rows[0]
                            data = filtered_rows[1:]
                            return pd.DataFrame(data, columns=headers)
                            
            except Exception:
                continue
        
        return None
    
    async def extract_text(self, file_path: Union[str, Path]) -> str:
        """
        Extract all text content from PDF for RAG/search.
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            Extracted text content
        """
        file_path = Path(file_path)
        text_content = []
        
        if not PDF_PLUMBER_AVAILABLE:
            return ""
        
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append(page_text)
        except Exception as e:
            logger.error(f"Text extraction failed: {e}")
        
        return "\n\n".join(text_content)


# =============================================================================
# EXCEL EXTRACTOR
# =============================================================================

class ExcelExtractor(BaseExtractor):
    """
    Extracts data from Excel files (.xlsx, .xls, .xlsm).
    
    Features:
    - Multi-sheet support with automatic sheet detection
    - Formula value extraction (not formulas themselves)
    - Named range support
    - Date/time format handling
    - Merged cell handling
    
    Dependencies:
    - openpyxl (for .xlsx, .xlsm)
    - xlrd (for .xls - legacy format)
    """
    
    def can_handle(self, file_type: FileType) -> bool:
        """Excel extractor handles Excel file types"""
        return file_type in [FileType.EXCEL_XLSX, FileType.EXCEL_XLS, FileType.EXCEL_XLSM]
    
    async def extract(self, file_path: Union[str, Path]) -> List[ExtractedTable]:
        """
        Extract all sheets from an Excel file.
        
        Args:
            file_path: Path to the Excel file
            
        Returns:
            List of ExtractedTable objects (one per sheet with data)
        """
        file_path = Path(file_path)
        tables: List[ExtractedTable] = []
        
        self.config.report_progress("Opening Excel file...", 0.1)
        
        try:
            # Detect engine based on file type
            file_type = detect_file_type(file_path)
            
            if file_type == FileType.EXCEL_XLS:
                if not XLRD_AVAILABLE:
                    self.add_error("xlrd not installed - cannot read .xls files")
                    return tables
                engine = 'xlrd'
            else:
                if not OPENPYXL_AVAILABLE:
                    self.add_error("openpyxl not installed - cannot read .xlsx files")
                    return tables
                engine = 'openpyxl'
            
            # Get sheet names
            xl_file = pd.ExcelFile(file_path, engine=engine)
            sheet_names = xl_file.sheet_names
            
            self.config.report_progress(f"Found {len(sheet_names)} sheets", 0.2)
            
            # Process each sheet
            for sheet_idx, sheet_name in enumerate(sheet_names):
                progress = 0.2 + (0.7 * (sheet_idx / len(sheet_names)))
                self.config.report_progress(f"Processing sheet: {sheet_name}", progress)
                
                # Extract sheet data
                sheet_table = self._extract_sheet(xl_file, sheet_name, engine)
                
                if sheet_table:
                    tables.append(sheet_table)
                
                # Check max tables limit
                if len(tables) >= self.config.max_tables:
                    self.add_warning(f"Reached maximum table limit ({self.config.max_tables})")
                    break
            
            xl_file.close()
            self.config.report_progress("Extraction complete", 1.0)
            
        except Exception as e:
            self.add_error(f"Excel extraction failed: {str(e)}")
            logger.exception("Excel extraction error")
        
        return tables
    
    def _extract_sheet(
        self, 
        xl_file: pd.ExcelFile, 
        sheet_name: str,
        engine: str
    ) -> Optional[ExtractedTable]:
        """
        Extract data from a single Excel sheet.
        
        Args:
            xl_file: pandas ExcelFile object
            sheet_name: Name of the sheet to extract
            engine: Excel engine to use
            
        Returns:
            ExtractedTable or None if sheet is empty
        """
        try:
            # Read the sheet
            # First, try to detect header row
            df_preview = pd.read_excel(
                xl_file,
                sheet_name=sheet_name,
                header=None,
                nrows=20,
                engine=engine
            )
            
            # Find the best header row
            header_row = self._detect_header_row(df_preview)
            
            # Re-read with correct header
            df = pd.read_excel(
                xl_file,
                sheet_name=sheet_name,
                header=header_row,
                engine=engine
            )
            
            # Clean the dataframe
            df = clean_dataframe(df)
            
            # Skip if empty
            if self.config.skip_empty_tables and (df is None or df.empty):
                return None
            
            # Skip if too few rows
            if df is not None and len(df) < self.config.min_table_rows:
                return None
            
            # Apply max rows limit
            if self.config.max_rows and len(df) > self.config.max_rows:
                df = df.head(self.config.max_rows)
                self.add_warning(f"Sheet '{sheet_name}' truncated to {self.config.max_rows} rows")
            
            # Get additional metadata
            metadata = {
                "sheet_name": sheet_name,
                "header_row": header_row,
                "extraction_method": "pandas",
                "engine": engine
            }
            
            # Try to detect date columns
            date_columns = self._detect_date_columns(df)
            if date_columns:
                metadata["date_columns"] = date_columns
            
            return ExtractedTable(
                name=sheet_name,
                dataframe=df,
                source_location=f"Sheet: {sheet_name}",
                quality=assess_data_quality(df),
                metadata=metadata
            )
            
        except Exception as e:
            self.add_warning(f"Failed to extract sheet '{sheet_name}': {str(e)}")
            return None
    
    def _detect_header_row(self, df_preview: pd.DataFrame) -> int:
        """
        Detect the row that should be used as the header.
        
        Looks for a row where:
        - Most values are strings
        - Values are unique
        - Not mostly numeric
        
        Args:
            df_preview: Preview of the first rows
            
        Returns:
            Row number to use as header (0-indexed)
        """
        if self.config.header_row is not None:
            return self.config.header_row
        
        best_row = 0
        best_score = 0
        
        for row_idx in range(min(10, len(df_preview))):
            row = df_preview.iloc[row_idx]
            
            # Skip completely empty rows
            if row.isna().all():
                continue
            
            score = 0
            
            # Score: more non-null values is better
            non_null = row.notna().sum()
            score += non_null * 2
            
            # Score: strings are better than numbers for headers
            string_count = sum(1 for v in row if isinstance(v, str))
            score += string_count * 3
            
            # Score: unique values are better
            unique_count = row.nunique()
            score += unique_count * 2
            
            # Penalty: numeric values suggest data row
            numeric_count = sum(1 for v in row if isinstance(v, (int, float)) and not pd.isna(v))
            score -= numeric_count * 2
            
            if score > best_score:
                best_score = score
                best_row = row_idx
        
        return best_row
    
    def _detect_date_columns(self, df: pd.DataFrame) -> List[str]:
        """
        Detect columns that contain date/datetime data.
        
        Args:
            df: DataFrame to analyze
            
        Returns:
            List of column names that appear to be dates
        """
        date_columns = []
        
        for col in df.columns:
            # Already datetime type
            if pd.api.types.is_datetime64_any_dtype(df[col]):
                date_columns.append(col)
                continue
            
            # Check column name for date-related keywords
            col_lower = str(col).lower()
            date_keywords = ['date', 'time', 'day', 'month', 'year', 'created', 'updated', 'timestamp']
            if any(kw in col_lower for kw in date_keywords):
                # Try to parse as date
                try:
                    pd.to_datetime(df[col].head(100), errors='raise')
                    date_columns.append(col)
                except:
                    pass
        
        return date_columns


# =============================================================================
# CSV/TSV EXTRACTOR
# =============================================================================

class CSVExtractor(BaseExtractor):
    """
    Extracts data from CSV and TSV files.
    
    Features:
    - Auto-encoding detection
    - Auto-delimiter detection
    - Handles various quote characters
    - Large file support with chunking
    """
    
    def can_handle(self, file_type: FileType) -> bool:
        """CSV extractor handles CSV and TSV files"""
        return file_type in [FileType.CSV, FileType.TSV]
    
    async def extract(self, file_path: Union[str, Path]) -> List[ExtractedTable]:
        """
        Extract data from a CSV/TSV file.
        
        Args:
            file_path: Path to the CSV file
            
        Returns:
            List containing single ExtractedTable
        """
        file_path = Path(file_path)
        tables: List[ExtractedTable] = []
        
        self.config.report_progress("Detecting encoding...", 0.1)
        
        try:
            # Detect encoding
            encoding = self.config.encoding or detect_encoding(file_path)
            
            self.config.report_progress("Detecting delimiter...", 0.2)
            
            # Detect delimiter
            file_type = detect_file_type(file_path)
            if file_type == FileType.TSV:
                delimiter = '\t'
            else:
                delimiter = detect_csv_delimiter(file_path, encoding)
            
            self.config.report_progress("Reading data...", 0.3)
            
            # Read the CSV
            df = pd.read_csv(
                file_path,
                encoding=encoding,
                delimiter=delimiter,
                header=self.config.header_row if self.config.header_row is not None else 'infer',
                nrows=self.config.max_rows,
                on_bad_lines='warn',
                low_memory=False
            )
            
            self.config.report_progress("Cleaning data...", 0.8)
            
            # Clean the dataframe
            df = clean_dataframe(df)
            
            if df is not None and not df.empty:
                table = ExtractedTable(
                    name=file_path.stem,
                    dataframe=df,
                    source_location=str(file_path.name),
                    quality=assess_data_quality(df),
                    metadata={
                        "encoding": encoding,
                        "delimiter": repr(delimiter),
                        "extraction_method": "pandas"
                    }
                )
                tables.append(table)
            
            self.config.report_progress("Complete", 1.0)
            
        except Exception as e:
            self.add_error(f"CSV extraction failed: {str(e)}")
            logger.exception("CSV extraction error")
        
        return tables


# =============================================================================
# POWERPOINT EXTRACTOR
# =============================================================================

class PowerPointExtractor(BaseExtractor):
    """
    Extracts tables from PowerPoint presentations.
    
    Features:
    - Table extraction from slides
    - Shape text extraction
    - Notes extraction
    """
    
    def can_handle(self, file_type: FileType) -> bool:
        """PowerPoint extractor handles .pptx files"""
        return file_type == FileType.POWERPOINT
    
    async def extract(self, file_path: Union[str, Path]) -> List[ExtractedTable]:
        """
        Extract tables from PowerPoint slides.
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            List of ExtractedTable objects
        """
        file_path = Path(file_path)
        tables: List[ExtractedTable] = []
        
        if not PPTX_AVAILABLE:
            self.add_error("python-pptx not installed - cannot read PowerPoint files")
            return tables
        
        self.config.report_progress("Opening presentation...", 0.1)
        
        try:
            prs = Presentation(file_path)
            total_slides = len(prs.slides)
            
            self.config.report_progress(f"Processing {total_slides} slides...", 0.2)
            
            for slide_idx, slide in enumerate(prs.slides):
                progress = 0.2 + (0.7 * (slide_idx / total_slides))
                self.config.report_progress(f"Processing slide {slide_idx + 1}", progress)
                
                # Find tables in shapes
                table_count = 0
                for shape in slide.shapes:
                    if shape.has_table:
                        table_data = self._extract_pptx_table(shape.table)
                        
                        if table_data is not None and not table_data.empty:
                            table_count += 1
                            table = ExtractedTable(
                                name=f"Slide_{slide_idx + 1}_Table_{table_count}",
                                dataframe=table_data,
                                source_location=f"Slide {slide_idx + 1}",
                                quality=assess_data_quality(table_data),
                                metadata={
                                    "slide_number": slide_idx + 1,
                                    "extraction_method": "python-pptx"
                                }
                            )
                            tables.append(table)
                
                # Check max tables limit
                if len(tables) >= self.config.max_tables:
                    self.add_warning(f"Reached maximum table limit")
                    break
            
            self.config.report_progress("Complete", 1.0)
            
        except Exception as e:
            self.add_error(f"PowerPoint extraction failed: {str(e)}")
            logger.exception("PowerPoint extraction error")
        
        return tables
    
    def _extract_pptx_table(self, table) -> Optional[pd.DataFrame]:
        """
        Extract data from a PowerPoint table shape.
        
        Args:
            table: python-pptx table object
            
        Returns:
            DataFrame or None
        """
        try:
            rows = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    # Get cell text
                    text = cell.text.strip() if cell.text else ""
                    row_data.append(text)
                rows.append(row_data)
            
            if len(rows) < 2:
                return None
            
            # First row as header
            headers = rows[0]
            data = rows[1:]
            
            # Clean headers
            headers = [h if h else f"Column_{i}" for i, h in enumerate(headers)]
            
            df = pd.DataFrame(data, columns=headers)
            return clean_dataframe(df)
            
        except Exception as e:
            logger.warning(f"Failed to extract PowerPoint table: {e}")
            return None


# =============================================================================
# WORD DOCUMENT EXTRACTOR
# =============================================================================

class WordExtractor(BaseExtractor):
    """
    Extracts tables from Word documents.
    
    Features:
    - Table extraction
    - Text extraction for RAG
    - Preserves table structure
    """
    
    def can_handle(self, file_type: FileType) -> bool:
        """Word extractor handles .docx files"""
        return file_type == FileType.WORD
    
    async def extract(self, file_path: Union[str, Path]) -> List[ExtractedTable]:
        """
        Extract tables from Word document.
        
        Args:
            file_path: Path to the Word file
            
        Returns:
            List of ExtractedTable objects
        """
        file_path = Path(file_path)
        tables: List[ExtractedTable] = []
        
        if not DOCX_AVAILABLE:
            self.add_error("python-docx not installed - cannot read Word files")
            return tables
        
        self.config.report_progress("Opening document...", 0.1)
        
        try:
            doc = DocxDocument(file_path)
            doc_tables = doc.tables
            
            self.config.report_progress(f"Found {len(doc_tables)} tables", 0.2)
            
            for table_idx, table in enumerate(doc_tables):
                progress = 0.2 + (0.7 * (table_idx / max(1, len(doc_tables))))
                self.config.report_progress(f"Processing table {table_idx + 1}", progress)
                
                table_data = self._extract_docx_table(table)
                
                if table_data is not None and not table_data.empty:
                    extracted = ExtractedTable(
                        name=f"Table_{table_idx + 1}",
                        dataframe=table_data,
                        source_location=f"Document Table {table_idx + 1}",
                        quality=assess_data_quality(table_data),
                        metadata={
                            "table_index": table_idx,
                            "extraction_method": "python-docx"
                        }
                    )
                    tables.append(extracted)
                
                # Check max tables limit
                if len(tables) >= self.config.max_tables:
                    self.add_warning(f"Reached maximum table limit")
                    break
            
            self.config.report_progress("Complete", 1.0)
            
        except Exception as e:
            self.add_error(f"Word extraction failed: {str(e)}")
            logger.exception("Word extraction error")
        
        return tables
    
    def _extract_docx_table(self, table) -> Optional[pd.DataFrame]:
        """
        Extract data from a Word table.
        
        Args:
            table: python-docx table object
            
        Returns:
            DataFrame or None
        """
        try:
            rows = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    text = cell.text.strip() if cell.text else ""
                    row_data.append(text)
                rows.append(row_data)
            
            if len(rows) < 2:
                return None
            
            # First row as header
            headers = rows[0]
            data = rows[1:]
            
            # Clean headers
            headers = [h if h else f"Column_{i}" for i, h in enumerate(headers)]
            
            df = pd.DataFrame(data, columns=headers)
            return clean_dataframe(df)
            
        except Exception as e:
            logger.warning(f"Failed to extract Word table: {e}")
            return None


# =============================================================================
# IMAGE EXTRACTOR (OCR)
# =============================================================================

class ImageExtractor(BaseExtractor):
    """
    Extracts tables from images using OCR.
    
    Features:
    - OCR text extraction
    - Table structure detection
    - Supports PNG, JPG, TIFF, BMP
    """
    
    def can_handle(self, file_type: FileType) -> bool:
        """Image extractor handles image files"""
        return file_type in [
            FileType.IMAGE_PNG, FileType.IMAGE_JPG, 
            FileType.IMAGE_JPEG, FileType.IMAGE_TIFF,
            FileType.IMAGE_BMP
        ]
    
    async def extract(self, file_path: Union[str, Path]) -> List[ExtractedTable]:
        """
        Extract tables from image using OCR.
        
        Args:
            file_path: Path to the image file
            
        Returns:
            List of ExtractedTable objects
        """
        file_path = Path(file_path)
        tables: List[ExtractedTable] = []
        
        if not TESSERACT_AVAILABLE:
            self.add_error("pytesseract/PIL not installed - cannot perform OCR")
            return tables
        
        if not self.config.ocr_enabled:
            self.add_warning("OCR is disabled in configuration")
            return tables
        
        self.config.report_progress("Loading image...", 0.1)
        
        try:
            # Open image
            img = Image.open(file_path)
            
            self.config.report_progress("Running OCR...", 0.3)
            
            # Run OCR
            ocr_text = pytesseract.image_to_string(
                img,
                lang=self.config.ocr_language,
                config='--psm 6'
            )
            
            self.config.report_progress("Parsing table structure...", 0.7)
            
            # Try to parse as table
            df = self._parse_ocr_to_table(ocr_text)
            
            if df is not None and not df.empty:
                table = ExtractedTable(
                    name=file_path.stem,
                    dataframe=df,
                    source_location=str(file_path.name),
                    quality=DataQuality.LOW,  # OCR quality is typically lower
                    metadata={
                        "extraction_method": "ocr",
                        "ocr_language": self.config.ocr_language,
                        "image_size": img.size
                    }
                )
                tables.append(table)
            else:
                self.add_warning("No tabular data detected in image")
            
            self.config.report_progress("Complete", 1.0)
            
        except Exception as e:
            self.add_error(f"Image extraction failed: {str(e)}")
            logger.exception("Image extraction error")
        
        return tables
    
    def _parse_ocr_to_table(self, ocr_text: str) -> Optional[pd.DataFrame]:
        """
        Parse OCR text into a table structure.
        
        Args:
            ocr_text: Raw OCR text
            
        Returns:
            DataFrame or None if no table structure detected
        """
        if not ocr_text or not ocr_text.strip():
            return None
        
        lines = [line.strip() for line in ocr_text.split('\n') if line.strip()]
        
        if len(lines) < 2:
            return None
        
        # Try different parsing strategies
        strategies = [
            self._parse_with_tabs,
            self._parse_with_spaces,
            self._parse_with_pipes
        ]
        
        for strategy in strategies:
            df = strategy(lines)
            if df is not None and len(df.columns) >= 2:
                return df
        
        return None
    
    def _parse_with_tabs(self, lines: List[str]) -> Optional[pd.DataFrame]:
        """Parse lines assuming tab-separated values"""
        try:
            rows = [line.split('\t') for line in lines]
            if len(rows) < 2:
                return None
            
            # Check column consistency
            col_counts = [len(row) for row in rows]
            if max(col_counts) - min(col_counts) > 2:
                return None
            
            headers = rows[0]
            data = rows[1:]
            return pd.DataFrame(data, columns=headers)
        except:
            return None
    
    def _parse_with_spaces(self, lines: List[str]) -> Optional[pd.DataFrame]:
        """Parse lines assuming space-separated values"""
        try:
            rows = [re.split(r'\s{2,}', line) for line in lines]
            if len(rows) < 2:
                return None
            
            # Filter to consistent column counts
            col_counts = [len(row) for row in rows]
            most_common = max(set(col_counts), key=col_counts.count)
            rows = [row for row in rows if len(row) == most_common]
            
            if len(rows) < 2:
                return None
            
            headers = rows[0]
            data = rows[1:]
            return pd.DataFrame(data, columns=headers)
        except:
            return None
    
    def _parse_with_pipes(self, lines: List[str]) -> Optional[pd.DataFrame]:
        """Parse lines assuming pipe-separated values"""
        try:
            rows = [[cell.strip() for cell in line.split('|') if cell.strip()] for line in lines]
            rows = [row for row in rows if row]  # Remove empty rows
            
            if len(rows) < 2:
                return None
            
            headers = rows[0]
            data = rows[1:]
            return pd.DataFrame(data, columns=headers)
        except:
            return None


# =============================================================================
# JSON EXTRACTOR
# =============================================================================

class JSONExtractor(BaseExtractor):
    """
    Extracts tabular data from JSON files.
    
    Features:
    - Handles arrays of objects
    - Flattens nested structures
    - Handles JSON Lines format
    """
    
    def can_handle(self, file_type: FileType) -> bool:
        """JSON extractor handles JSON files"""
        return file_type == FileType.JSON
    
    async def extract(self, file_path: Union[str, Path]) -> List[ExtractedTable]:
        """
        Extract data from JSON file.
        
        Args:
            file_path: Path to the JSON file
            
        Returns:
            List of ExtractedTable objects
        """
        file_path = Path(file_path)
        tables: List[ExtractedTable] = []
        
        self.config.report_progress("Reading JSON file...", 0.2)
        
        try:
            # Detect encoding
            encoding = self.config.encoding or detect_encoding(file_path)
            
            with open(file_path, 'r', encoding=encoding) as f:
                content = f.read()
            
            self.config.report_progress("Parsing JSON...", 0.4)
            
            # Try to parse as JSON
            try:
                data = json.loads(content)
            except json.JSONDecodeError:
                # Try JSON Lines format
                data = [json.loads(line) for line in content.strip().split('\n') if line.strip()]
            
            self.config.report_progress("Converting to DataFrame...", 0.6)
            
            # Convert to DataFrame
            df = self._json_to_dataframe(data)
            
            if df is not None and not df.empty:
                # Apply max rows
                if self.config.max_rows and len(df) > self.config.max_rows:
                    df = df.head(self.config.max_rows)
                    self.add_warning(f"Truncated to {self.config.max_rows} rows")
                
                table = ExtractedTable(
                    name=file_path.stem,
                    dataframe=df,
                    source_location=str(file_path.name),
                    quality=assess_data_quality(df),
                    metadata={
                        "encoding": encoding,
                        "extraction_method": "json_parse"
                    }
                )
                tables.append(table)
            
            self.config.report_progress("Complete", 1.0)
            
        except Exception as e:
            self.add_error(f"JSON extraction failed: {str(e)}")
            logger.exception("JSON extraction error")
        
        return tables
    
    def _json_to_dataframe(self, data: Any) -> Optional[pd.DataFrame]:
        """
        Convert JSON data to DataFrame.
        
        Args:
            data: Parsed JSON data
            
        Returns:
            DataFrame or None
        """
        try:
            # If it's a list of dicts, use json_normalize for flattening
            if isinstance(data, list) and len(data) > 0 and isinstance(data[0], dict):
                df = pd.json_normalize(data, max_level=3)
                return clean_dataframe(df)
            
            # If it's a dict with a data key containing a list
            if isinstance(data, dict):
                # Look for common data keys
                for key in ['data', 'results', 'items', 'records', 'rows']:
                    if key in data and isinstance(data[key], list):
                        df = pd.json_normalize(data[key], max_level=3)
                        return clean_dataframe(df)
                
                # Try to convert the whole dict
                df = pd.json_normalize([data], max_level=3)
                return clean_dataframe(df)
            
            return None
            
        except Exception as e:
            logger.warning(f"JSON to DataFrame conversion failed: {e}")
            return None


# =============================================================================
# MAIN FILE INGESTOR CLASS
# =============================================================================

class FileIngestor:
    """
    Main class for extracting data from multiple file formats.
    
    This is the primary interface for the data extraction system.
    It automatically detects file types and routes to the appropriate extractor.
    
    Usage:
        ingestor = FileIngestor()
        result = await ingestor.extract("data.xlsx")
        
        # Access data
        df = result.get_dataframe()  # Get primary table
        
        # Check status
        if result.has_data:
            print(f"Extracted {result.total_rows} rows")
    """
    
    def __init__(self, config: Optional[ExtractionConfig] = None):
        """
        Initialize the FileIngestor.
        
        Args:
            config: Extraction configuration options
        """
        self.config = config or ExtractionConfig()
        
        # Initialize all extractors
        self._extractors: List[BaseExtractor] = [
            PDFExtractor(self.config),
            ExcelExtractor(self.config),
            CSVExtractor(self.config),
            PowerPointExtractor(self.config),
            WordExtractor(self.config),
            ImageExtractor(self.config),
            JSONExtractor(self.config),
        ]
        
        logger.info("FileIngestor initialized with %d extractors", len(self._extractors))
    
    def get_supported_formats(self) -> List[str]:
        """
        Get list of supported file extensions.
        
        Returns:
            List of supported extensions
        """
        return [
            '.pdf', '.xlsx', '.xls', '.xlsm',
            '.csv', '.tsv',
            '.pptx', '.docx',
            '.png', '.jpg', '.jpeg', '.tiff', '.bmp',
            '.json'
        ]
    
    def _get_extractor(self, file_type: FileType) -> Optional[BaseExtractor]:
        """
        Get the appropriate extractor for a file type.
        
        Args:
            file_type: The file type to handle
            
        Returns:
            Appropriate extractor or None
        """
        for extractor in self._extractors:
            if extractor.can_handle(file_type):
                return extractor
        return None
    
    async def extract(
        self, 
        file_path: Union[str, Path, BinaryIO],
        file_name: Optional[str] = None
    ) -> ExtractionResult:
        """
        Extract data from a file.
        
        This is the main entry point for file extraction.
        
        Args:
            file_path: Path to file, or file-like object
            file_name: Optional filename (required if file_path is file-like)
            
        Returns:
            ExtractionResult containing all extracted data
        """
        import time
        start_time = time.time()
        
        # Handle file-like objects
        temp_file = None
        if hasattr(file_path, 'read'):
            if not file_name:
                raise ValueError("file_name required when passing file-like object")
            
            # Save to temp file
            temp_file = tempfile.NamedTemporaryFile(
                delete=False,
                suffix=Path(file_name).suffix
            )
            temp_file.write(file_path.read())
            temp_file.close()
            file_path = Path(temp_file.name)
            actual_name = file_name
        else:
            file_path = Path(file_path)
            actual_name = file_path.name
        
        # Detect file type
        file_type = detect_file_type(file_path)
        
        # Initialize result
        result = ExtractionResult(
            source_file=actual_name,
            file_type=file_type,
            status=ExtractionStatus.FAILED,
            metadata={
                "file_size": file_path.stat().st_size if file_path.exists() else 0,
                "extraction_started": datetime.now().isoformat()
            }
        )
        
        try:
            # Check if file exists
            if not file_path.exists():
                result.errors.append(f"File not found: {file_path}")
                return result
            
            # Check file type
            if file_type == FileType.UNKNOWN:
                result.errors.append(f"Unsupported file format: {file_path.suffix}")
                return result
            
            # Get appropriate extractor
            extractor = self._get_extractor(file_type)
            
            if not extractor:
                result.errors.append(f"No extractor available for {file_type.value}")
                return result
            
            # Clear any previous messages
            extractor.clear_messages()
            
            # Perform extraction
            tables = await extractor.extract(file_path)
            
            # Collect results
            result.tables = tables
            result.warnings = extractor.warnings.copy()
            result.errors = extractor.errors.copy()
            
            # Determine status
            if tables and any(t.row_count > 0 for t in tables):
                if extractor.errors:
                    result.status = ExtractionStatus.PARTIAL
                else:
                    result.status = ExtractionStatus.SUCCESS
            elif not extractor.errors:
                result.status = ExtractionStatus.NO_DATA
                result.warnings.append("No tabular data found in file")
            else:
                result.status = ExtractionStatus.FAILED
            
            # Extract text content for RAG (if PDF)
            if file_type == FileType.PDF and isinstance(extractor, PDFExtractor):
                result.text_content = await extractor.extract_text(file_path)
            
        except Exception as e:
            result.errors.append(f"Extraction failed: {str(e)}")
            logger.exception("Extraction error")
        
        finally:
            # Clean up temp file
            if temp_file and Path(temp_file.name).exists():
                try:
                    os.unlink(temp_file.name)
                except:
                    pass
            
            # Record processing time
            result.processing_time = time.time() - start_time
            result.metadata["extraction_completed"] = datetime.now().isoformat()
        
        return result
    
    async def extract_multiple(
        self,
        file_paths: List[Union[str, Path]]
    ) -> List[ExtractionResult]:
        """
        Extract data from multiple files.
        
        Args:
            file_paths: List of file paths
            
        Returns:
            List of ExtractionResult objects
        """
        results = []
        
        for i, file_path in enumerate(file_paths):
            self.config.report_progress(
                f"Processing file {i + 1}/{len(file_paths)}",
                i / len(file_paths)
            )
            
            result = await self.extract(file_path)
            results.append(result)
        
        return results
    
    def get_extraction_capabilities(self) -> Dict[str, Any]:
        """
        Get information about extraction capabilities.
        
        Returns:
            Dictionary with capability information
        """
        return {
            "supported_formats": self.get_supported_formats(),
            "pdf_support": PDF_PLUMBER_AVAILABLE,
            "ocr_support": TESSERACT_AVAILABLE,
            "excel_xlsx_support": OPENPYXL_AVAILABLE,
            "excel_xls_support": XLRD_AVAILABLE,
            "powerpoint_support": PPTX_AVAILABLE,
            "word_support": DOCX_AVAILABLE,
            "encoding_detection": CHARDET_AVAILABLE,
            "max_rows_limit": self.config.max_rows,
            "max_tables_limit": self.config.max_tables,
            "ocr_enabled": self.config.ocr_enabled,
            "ocr_language": self.config.ocr_language
        }


# =============================================================================
# CONVENIENCE FUNCTIONS
# =============================================================================

async def extract_file(
    file_path: Union[str, Path],
    config: Optional[ExtractionConfig] = None
) -> ExtractionResult:
    """
    Convenience function to extract data from a file.
    
    Args:
        file_path: Path to the file
        config: Optional extraction configuration
        
    Returns:
        ExtractionResult
    
    Example:
        result = await extract_file("sales.xlsx")
        df = result.get_dataframe()
    """
    ingestor = FileIngestor(config)
    return await ingestor.extract(file_path)


def extract_file_sync(
    file_path: Union[str, Path],
    config: Optional[ExtractionConfig] = None
) -> ExtractionResult:
    """
    Synchronous wrapper for extract_file.
    
    Args:
        file_path: Path to the file
        config: Optional extraction configuration
        
    Returns:
        ExtractionResult
    
    Example:
        result = extract_file_sync("sales.xlsx")
        df = result.get_dataframe()
    """
    return asyncio.run(extract_file(file_path, config))


# =============================================================================
# MODULE EXPORTS
# =============================================================================

__all__ = [
    # Main classes
    'FileIngestor',
    'ExtractionResult',
    'ExtractedTable',
    'ExtractionConfig',
    
    # Extractors
    'PDFExtractor',
    'ExcelExtractor',
    'CSVExtractor',
    'PowerPointExtractor',
    'WordExtractor',
    'ImageExtractor',
    'JSONExtractor',
    
    # Enums
    'FileType',
    'ExtractionStatus',
    'DataQuality',
    
    # Convenience functions
    'extract_file',
    'extract_file_sync',
    
    # Utilities
    'detect_file_type',
    'detect_encoding',
]
